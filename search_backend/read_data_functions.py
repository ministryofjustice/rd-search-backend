"""
A collection of functions to help read data from PDFs, Word docs, and Powerpoints.
"""

import re
from io import BytesIO

from docx import Document
from pdfminer.converter import TextConverter
from pdfminer.layout import LAParams
from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.pdfpage import PDFPage
from pptx import Presentation

from search_backend.s3client import S3Client


def _read_pdf_gen(f, title, fname):

    """
    Generator to read pages from PDF
    """

    # Define generator from which pages can be read
    page_gen = PDFPage.get_pages(f)

    # This dictionary specifies strings that need replacing in the extracted text.
    # This is because the PDF reader can't interpret certain punctuation marks.
    replace_dict = {
        '\x0c': '',
        # '\n': ' ',
        '\xe2\x80\x93': '-',
        '\xc2\xa3': '£',
        '\uf0b7': '•',
        '': '',
    }

    # Iterate over the pdf pages
    for pageNumber, page in enumerate(page_gen):

        # Define some objects required by the pdf reader
        laparams = LAParams()
        rsrcmgr = PDFResourceManager()
        retstr = BytesIO()
        device = TextConverter(rsrcmgr, retstr, codec='utf-8', laparams=laparams)
        interpreter = PDFPageInterpreter(rsrcmgr, device)

        # Extract the text
        interpreter.process_page(page)
        data = retstr.getvalue()
        body = data.decode('utf-8')

        # Replace weird sets of characters that get introduced by the pdf reader
        # when it can't interpret a punctuation mark
        for match, repl in replace_dict.items():
            body = body.replace(match, repl)

        # Replace multiple spaces with single space
        body = re.sub(r'\s+', ' ', body.strip())

        if body == '':
            # Skip pages with no text
            continue
        if body is None:
            continue
        elif bool(re.search('^contents', body.lower())) & (pageNumber <= 5):
            # Skip contents pages
            continue
        elif (pageNumber == 0) & (len(body.split(' ')) <= 8):
            # Skip what is likely to be a title page
            continue
        else:
            # Add page info to a dataframe
            page_dict = {
                'meta': {
                    'title': title,
                    'path': fname,
                    'page': pageNumber+1,
                },
                'content': body,
            }

        retstr.close()
        device.close()

        yield page_dict


def _read_word(f, title, fname):
    
    """
    Generator to read paragraphs from Word doc.

    Currently combining all text for across the doc together, as Word format doesn't
    break down into pages, and breaking into paragraphs means individual bullet points
    get separated.
    """
    
    document = Document(f)
    word_text = []

    # Iterate over paragraphs
    for ii, para in enumerate(document.paragraphs):

        para_text = para.text.strip()
        para_text = para_text.replace('\xa0', '')
        if para_text == '':
            continue
        elif para.style.name in ['Title', 'Subtitle']:
            continue
        else:
            if para.style.name == 'List Paragraph':
                # Bullet point formatting
                para_text = ' - ' + para_text

            word_text.append(para_text)

    # Iterate over any tables
    for table in document.tables:
        for row in table.rows:
            table_text = ' - '.join([cell.text for cell in row.cells if cell.text != ''])
            if table_text != '':
                word_text.append(table_text)

    word_text = '\n'.join(word_text)
    # Add page info to a dataframe
    word_dict = {
        'meta': {
            'title': title,
            'path': fname,
            'page': 0,
        },
        'content': word_text,
    }

    return word_dict


def _read_ppt_gen(f, title, fname):
    
    """
    Generator to read paragraphs from Powerpoint doc
    """

    ppt = Presentation(f)

    for ii, slide in enumerate(ppt.slides):

        # Skip the first slide, since it's just a title slide
        if ii == 0:
            continue

        slide_text = []

        # Add all text from the slide to a list
        for shape in slide.shapes:

            if shape.has_text_frame:
                for jj, para in enumerate(shape.text_frame.paragraphs):
                    para_text = ' '.join([run.text for run in para.runs])
                    para_text = re.sub(r'\s+', ' ', para_text).strip()
                    para_text = para_text.replace('\xa0', '')
                    if para_text != '':
                        # Make sure the paragraph contains text
                        slide_text.append(para_text)

            # Check for a table, as text formatted via a table won't be captured by the above
            if shape.has_table:
                for row in shape.table.rows:
                    table_text = ' '.join([cell.text for cell in row.cells if cell.text != ''])
                    if table_text != '':
                        slide_text.append(table_text)

        # Combine slide text and yield if it meets requirements
        slide_text = '\n'.join(slide_text)
        if slide_text == '':
            # Skip slides with no text
            continue
        elif re.search('^contents|\ncontents$', slide_text.lower()):
            # Skip contents page slides
            continue
        elif len(slide_text.split(' ')) <= 8:
            # Skip what is likely to be a title slide
            continue
        else:
            # Add page info to a dataframe
            slide_dict = {
                'meta': {
                    'title': title,
                    'path': fname,
                    'page': ii+1,
                },
                'content': slide_text,
            }

            yield slide_dict
    

def read_docs(s3client: S3Client, fnames: list[str]):
    """
    Iterates through a list of docs and calls the appropriate function
    to extract the text from each page (if applicable) of each doc.

    Arguments:
     - s3client: S3 client set up with authentication
     - fnames: a list of keys in the bucket

    Returns:
    A list of dictionaries containing document text and metadata.
    """

    print("Reading documents...")

    data = []

    for fname in fnames:
        # Get the title from the filename
        title = fname.split('/')[-1]

        # Read in the pdf
        obj, _ = s3client.get_object(fname, prepend_prefix=False)

        fs = obj["Body"].read()

        with BytesIO(fs) as f:
            if re.search('.pdf$', fname):
                doc_list = [page for page in _read_pdf_gen(f, title, fname)]
            elif re.search('.doc$|.docx$', fname):
                doc_list = [_read_word(f, title, fname)]
            elif re.search('.ppt$|.pptx$', fname):
                doc_list = [para for para in _read_ppt_gen(f, title, fname)]
            else:
                print(f"File format not accepted for {fname}")

        data.append(doc_list)

    # Unpack so we don't have nested lists
    data = [page for doc in data for page in doc]

    return data
